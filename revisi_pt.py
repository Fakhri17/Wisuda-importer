import os
import re
import json
import pandas as pd
import random
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class GraduationPPTGenerator:
    DPI = 96  # konsisten dengan PowerPoint
    
    # POSISI FOTO - TENGAH FRAME MERAH
    # Sesuaikan nilai ini agar foto berada di tengah frame merah
    PHOTO_FRAME_W_CM = 5.0     # Lebar frame foto
    PHOTO_FRAME_H_CM = 7.0      # Tinggi frame foto
    FRAME_LEFT_CM = 7.0         # Posisi horizontal (tengah slide)
    FRAME_TOP_CM = 4.85          # Posisi vertikal (tengah frame merah)

    def __init__(self):
        self.templates = {
            'Non Predikat': 'templates/template-pt/Slide1.PNG',
            'CUMLAUDE': 'templates/template-pt/Slide2.PNG',
            'SUMMA CUMLAUDE': 'templates/template-pt/Slide3.PNG',
        }

    # =========================
    # Helpers ukuran & gambar
    # =========================
    def _set_slide_size_to_image_exact(self, prs, image_path, dpi=None):
        """Sesuaikan ukuran slide PERSIS dengan ukuran gambar (pixel -> inch @DPI)."""
        if dpi is None:
            dpi = self.DPI
        if not image_path or not os.path.exists(image_path):
            return
        with Image.open(image_path) as img:
            w_px, h_px = img.size
        prs.slide_width  = Inches(w_px / dpi)
        prs.slide_height = Inches(h_px / dpi)

    def _set_background_image(self, slide, image_path):
        """Pasang gambar sebagai latar dengan menambahkan picture di (0,0) ukuran native."""
        if not image_path or not os.path.exists(image_path):
            return
        try:
            slide.shapes.add_picture(image_path, 0, 0)
        except Exception as e:
            print(f"Error setting background image {image_path}: {e}")

    def _add_picture_fit(self, slide, image_path, left, top, frame_width, frame_height):
        """Tambahkan gambar agar pas di dalam frame tanpa distorsi (centered)."""
        try:
            with Image.open(image_path) as img:
                img_w, img_h = img.size
        except Exception as e:
            print(f"Error opening image {image_path}: {e}")
            return None

        frame_ratio = float(frame_width) / float(frame_height) if frame_height else 1.0
        img_ratio = img_w / img_h if img_h else 1.0

        if img_ratio > frame_ratio:
            width = int(frame_width)
            height = int(width / img_ratio)
        else:
            height = int(frame_height)
            width = int(height * img_ratio)

        int_left = int(left)
        int_top = int(top)
        int_fw = int(frame_width)
        int_fh = int(frame_height)
        offset_left = int_left + (int_fw - int(width)) // 2
        offset_top = int_top + (int_fh - int(height)) // 2

        try:
            return slide.shapes.add_picture(image_path, int(offset_left), int(offset_top), width=int(width), height=int(height))
        except Exception as e:
            print(f"Error adding fitted picture {image_path}: {e}")
            return None

    # =========================
    # Data helpers
    # =========================
    def read_excel_data(self, file_path):
        """Read Excel file and return DataFrame."""
        try:
            df = pd.read_excel(file_path)
            print(f"Successfully read {len(df)} rows from {file_path}")
            print(f"Columns: {list(df.columns)}")
            return df
        except Exception as e:
            print(f"Error reading Excel file: {e}")
            return None

    def get_predikat_template(self, predikat):
        """Determine template based on predikat kelulusan."""
        if pd.isna(predikat) or predikat == '':
            return 'Non Predikat'
        predikat_str = str(predikat).lower().strip()
        if 'summa' in predikat_str and 'cumlaude' in predikat_str:
            return 'SUMMA CUMLAUDE'
        elif 'cumlaude' in predikat_str and 'summa' not in predikat_str:
            return 'CUMLAUDE'
        else:
            return 'Non Predikat'

    def find_student_photo(self, nim, program_folder):
        """Find student photo based on NIM in program folder."""
        photo_path = os.path.join('photos', program_folder, f"{nim}_graduation_1.jpg")
        return photo_path if os.path.exists(photo_path) else None

    def extract_seat_position(self, tempat_duduk):
        """Extract seat position for ordering (format '1.1.L')."""
        if pd.isna(tempat_duduk) or str(tempat_duduk).strip() == '':
            return (999, 999, 'Z')
        try:
            parts = str(tempat_duduk).split('.')
            if len(parts) >= 3:
                row = int(parts[0]); seat = int(parts[1]); side = parts[2].upper()
                return (row, seat, side)
        except Exception:
            pass
        return (999, 999, 'Z')

    # =========================
    # Slide builders
    # =========================
    def create_slide(self, prs, student_data, photo_path):
        """Create a single slide for a student."""
        predikat = self.get_predikat_template(student_data.get('PREDIKAT KELULUSAN', ''))
        template_path = self.templates.get(predikat, self.templates['Non Predikat'])

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background full-bleed
        if os.path.exists(template_path):
            self._set_background_image(slide, template_path)

        # Posisi dan ukuran frame foto dalam CM
        frame_left, frame_top = Cm(self.FRAME_LEFT_CM), Cm(self.FRAME_TOP_CM)
        frame_w, frame_h = Cm(self.PHOTO_FRAME_W_CM), Cm(self.PHOTO_FRAME_H_CM)

        # FOTO: fit ke dalam frame merah (tengah)
        if photo_path and os.path.exists(photo_path):
            try:
                self._add_picture_fit(slide, photo_path, frame_left, frame_top, frame_w, frame_h)
            except Exception as e:
                print(f"Error adding photo {photo_path}: {e}")

        # Teks info mahasiswa & dosen
        self.add_student_info(slide, student_data)

        return slide

    def _add_textbox(self, slide, text, left, top, width, height, font_size=18, bold=False, upper=True, alignment=PP_ALIGN.LEFT):
        """Utility untuk menambah textbox konsisten."""
        if not text or str(text).strip().lower() == 'nan':
            return
        content = str(text).upper() if upper else str(text)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = content
        p.alignment = alignment

        if len(p.runs) == 0:
            run = p.add_run()
            run.text = content
            p.text = ''

        for run in p.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

    def add_student_info(self, slide, student_data):
        """Add student information to slide - POSISI BARU SESUAI TEMPLATE."""
        program = student_data.get('PROGRAM STUDI', '')
        nama = student_data.get('NAMA MAHASISWA', '')
        nim = student_data.get('NIM', '')
        ipk = student_data.get('IPK', '')
        tak = student_data.get('SKOR TAK', '')
        dosen_wali = student_data.get('Nama Dosen Wali', '')
        perusahaan = "PT. Mencari Cinta SeJATI"

        dosen_pembimbing1 = student_data.get('Nama Dosen Pembimbing 1', '')
        dosen_pembimbing2 = student_data.get('Nama Dosen Pembimbing 2', '')
        pembimbing_names = []
        if pd.notna(dosen_pembimbing1) and str(dosen_pembimbing1).strip() != '':
            pembimbing_names.append(str(dosen_pembimbing1))
        if pd.notna(dosen_pembimbing2) and str(dosen_pembimbing2).strip() != '':
            pembimbing_names.append(str(dosen_pembimbing2))

        
        # PROGRAM STUDI : 
        self._add_textbox(slide, program, Cm(4.5), Cm(2.95), Cm(10), Cm(1), font_size=14, bold=True, alignment=PP_ALIGN.CENTER)

        # NAMA MAHASISWA : 
        self._add_textbox(slide, nama, Cm(0.2), Cm(14.2), Cm(19), Cm(1), font_size=19, bold=True, alignment=PP_ALIGN.CENTER)

        # NIM : 
        self._add_textbox(slide, nim, Cm(4.3), Cm(15.25), Cm(4), Cm(0.8), font_size=14, bold=True)

        # IPK : 
        self._add_textbox(slide, ipk, Cm(12.5), Cm(15.25), Cm(2), Cm(0.8), font_size=14, bold=True)
        
        # TAK : 
        self._add_textbox(slide, tak, Cm(15.6), Cm(15.25), Cm(2), Cm(0.8), font_size=14, bold=True)

        # PERUSAHAAN : 
        self._add_textbox(slide, perusahaan, Cm(5.7), Cm(16.17), Cm(10), Cm(0.8), font_size=12, bold=True)
        
        # DOSEN WALI : 
        self._add_textbox(slide, dosen_wali, Cm(5.7), Cm(16.94), Cm(12), Cm(0.8), font_size=12, bold=True)
        
        # DOSEN PEMBIMBING : 
        # Dosen Pembimbing Names
        if len(pembimbing_names) > 0:
            names_box = slide.shapes.add_textbox(Cm(5.7), Cm(17.78), Cm(12), Cm(1.5))
            names_tf = names_box.text_frame
            names_tf.clear()

            # Baris pertama
            p_name = names_tf.paragraphs[0]
            p_name.alignment = PP_ALIGN.LEFT
            run_n1 = p_name.add_run()
            run_n1.text = str(pembimbing_names[0]).upper()
            run_n1.font.name = 'Arial'
            run_n1.font.size = Pt(12)
            run_n1.font.bold = True
            run_n1.font.color.rgb = RGBColor(0, 0, 0)

            # Baris berikutnya
            for extra in pembimbing_names[1:]:
                p_more = names_tf.add_paragraph()
                p_more.alignment = PP_ALIGN.LEFT
                run_more = p_more.add_run()
                run_more.text = str(extra).upper()
                run_more.font.name = 'Arial'
                run_more.font.size = Pt(12)
                run_more.font.bold = True
                run_more.font.color.rgb = RGBColor(0, 0, 0)

        
    # =========================
    # Pipeline
    # =========================
    def extract_seat_side(self, tempat_duduk):
        """Extract seat side (L/R) from seat position."""
        if pd.isna(tempat_duduk) or str(tempat_duduk).strip() == '':
            return 'Z'
        try:
            parts = str(tempat_duduk).split('.')
            if len(parts) >= 3:
                return parts[2].upper()
        except Exception:
            pass
        return 'Z'

    def get_predikat_priority(self, predikat):
        """Get priority for sorting (1=summa, 2=cumlaude, 3=non-predikat)."""
        if pd.isna(predikat) or predikat == '':
            return 3
        predikat_str = str(predikat).lower().strip()
        if 'summa' in predikat_str and 'cumlaude' in predikat_str:
            return 1
        elif 'cumlaude' in predikat_str and 'summa' not in predikat_str:
            return 2
        else:
            return 3

    def generate_ppt_revisi(self, df, output_dir='output_revisi_pt', test_mode=False):
        """Generate PPT files separated by session, with summa students in separate folder."""
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if test_mode:
            # Test mode: generate single PPT with test data
            print(f"\nProcessing test data...")
            program_data = df.copy()
            
            prs = Presentation()
            
            # Use Cumlaude template for testing
            template_path = self.templates['CUMLAUDE']
            self._set_slide_size_to_image_exact(prs, template_path)
            
            # Add single test slide
            student = program_data.iloc[0]
            nim = student.get('NIM', '')
            photo_path = self.find_student_photo(nim, 'S1 Rekayasa Perangkat Lunak')
            if photo_path:
                print(f"  Adding test slide for {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
            else:
                print(f"  Warning: Photo not found for test data (NIM: {nim})")
            self.create_slide(prs, student, photo_path)
            
            # Save test file
            output_file = os.path.join(output_dir, "TEST_POSITION.pptx")
            prs.save(output_file)
            print(f"  Saved: {output_file} (1 test slide)")
            return

        # Group by session
        sessions = ['Pagi', 'Siang']
        sides = ['L', 'R']
        
        for session in sessions:
            # Filter by session
            session_data = df[df['SESI'] == session].copy()
            if len(session_data) == 0:
                print(f"\nNo data for session: {session}")
                continue
            
            # Session folder name
            session_folder_name = f"Wisuda {session}"
            session_output_dir = os.path.join(output_dir, session_folder_name)
            if not os.path.exists(session_output_dir):
                os.makedirs(session_output_dir)
            
            print(f"\nProcessing session: {session}")
            
            # 1. CREATE SUMMA FOLDER - All summa cumlaude students in one PPT
            summa_data = session_data.copy()
            summa_data['is_summa'] = summa_data['PREDIKAT KELULUSAN'].apply(
                lambda x: 'summa' in str(x).lower() and 'cumlaude' in str(x).lower() if pd.notna(x) else False
            )
            summa_students = summa_data[summa_data['is_summa']].copy()
            
            if len(summa_students) > 0:
                print(f"  Creating summa folder with {len(summa_students)} students")
                
                # Create summa folder
                summa_output_dir = os.path.join(session_output_dir, 'summa')
                if not os.path.exists(summa_output_dir):
                    os.makedirs(summa_output_dir)
                
                # Sort summa students by seat position
                summa_students['seat_sort'] = summa_students['TEMPAT DUDUK'].apply(self.extract_seat_position)
                summa_students = summa_students.sort_values('seat_sort')
                
                # Create summa PPT
                prs_summa = Presentation()
                try:
                    first_row = summa_students.iloc[0]
                    first_template_name = self.get_predikat_template(first_row.get('PREDIKAT KELULUSAN', ''))
                    first_template_path = self.templates.get(first_template_name, self.templates['SUMMA CUMLAUDE'])
                except Exception:
                    first_template_path = self.templates['SUMMA CUMLAUDE']
                self._set_slide_size_to_image_exact(prs_summa, first_template_path)
                
                # Add summa slides
                for _, student in summa_students.iterrows():
                    nim = student.get('NIM', '')
                    program = student.get('PROGRAM STUDI', '')
                    photo_path = self.find_student_photo(nim, program)
                    if photo_path:
                        print(f"    Adding summa slide for {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
                    else:
                        print(f"    Warning: Photo not found for summa student {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
                    self.create_slide(prs_summa, student, photo_path)
                
                # Save summa PPT
                summa_output_file = os.path.join(summa_output_dir, 'summa.pptx')
                prs_summa.save(summa_output_file)
                print(f"    Saved: {summa_output_file} ({len(summa_students)} slides)")
            else:
                print(f"  No summa cumlaude students found in {session} session")
            
            # 2. CREATE PROGRAM FOLDERS - Exclude summa students
            non_summa_data = session_data[~session_data['PREDIKAT KELULUSAN'].apply(
                lambda x: 'summa' in str(x).lower() and 'cumlaude' in str(x).lower() if pd.notna(x) else False
            )].copy()
            
            if len(non_summa_data) == 0:
                print(f"  No non-summa students found in {session} session")
                continue
            
            # Iterate programs in this session (excluding summa students)
            programs = [p for p in non_summa_data['PROGRAM STUDI'].dropna().unique() if str(p).strip() != '']
            print(f"  Processing {len(programs)} programs (excluding summa students)")
            
            for program in programs:
                prog_data = non_summa_data[non_summa_data['PROGRAM STUDI'] == program].copy()
                if len(prog_data) == 0:
                    continue
                
                # Prepare program folder
                safe_program_name = re.sub(r'[^\w\s-]', '', str(program)).strip()
                safe_program_name = re.sub(r'[-\s]+', '_', safe_program_name)
                program_output_dir = os.path.join(session_output_dir, safe_program_name)
                if not os.path.exists(program_output_dir):
                    os.makedirs(program_output_dir)
                
                print(f"    Program: {program} -> {len(prog_data)} students (non-summa)")
                
                # Produce two PPT files per program: L and R
                for side in sides:
                    side_data = prog_data.copy()
                    side_data['seat_side'] = side_data['TEMPAT DUDUK'].apply(self.extract_seat_side)
                    side_data = side_data[side_data['seat_side'] == side].copy()
                    if len(side_data) == 0:
                        print(f"      Skip Duduk {side}: no students")
                        continue
                    
                    # Sort by predikat then seat position
                    side_data['seat_sort'] = side_data['TEMPAT DUDUK'].apply(self.extract_seat_position)
                    side_data['predikat_priority'] = side_data['PREDIKAT KELULUSAN'].apply(self.get_predikat_priority)
                    side_data = side_data.sort_values(['predikat_priority', 'seat_sort'])
                    print(f"      Duduk {side}: {len(side_data)} students")
                    
                    # Create presentation with first student's template for exact slide size
                    prs = Presentation()
                    try:
                        first_row = side_data.iloc[0]
                        first_template_name = self.get_predikat_template(first_row.get('PREDIKAT KELULUSAN', ''))
                        first_template_path = self.templates.get(first_template_name, self.templates['Non Predikat'])
                    except Exception:
                        first_template_path = self.templates['Non Predikat']
                    self._set_slide_size_to_image_exact(prs, first_template_path)
                    
                    # Add slides
                    for _, student in side_data.iterrows():
                        nim = student.get('NIM', '')
                        photo_path = self.find_student_photo(nim, program)
                        if photo_path:
                            print(f"        Adding slide for {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
                        else:
                            print(f"        Warning: Photo not found for {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
                        self.create_slide(prs, student, photo_path)
                    
                    # Save to program folder
                    file_name = f"duduk_{side.lower()}.pptx"
                    output_file = os.path.join(program_output_dir, file_name)
                    prs.save(output_file)
                    print(f"      Saved: {output_file} ({len(side_data)} slides)")


    def create_test_data(self):
        """Create random test data for testing textbox positions."""
        test_data = {
            'PROGRAM STUDI': 'S1 Rekayasa Perangkat Lunak',
            'NAMA MAHASISWA': 'JOHN DOE TESTING DAN TESTING',
            'NIM': '1201200001',
            'IPK': '3.85',
            'SKOR TAK': '450',
            'Nama Dosen Wali': 'Dr. Ahmad Wijaya, S.T., M.T.',
            'Nama Dosen Pembimbing 1': 'Prof. Dr. Budi Santoso, S.T., M.T.',
            'Nama Dosen Pembimbing 2': 'Dr. Citra Dewi, S.T., M.Kom.',
            'PREDIKAT KELULUSAN': 'Cumlaude',
            'TEMPAT DUDUK': '1.1.L',
            'SESI': 'Pagi',
            'PERUSAHAAN': 'PT. Mencari Cinta Sejati'
        }
        return pd.DataFrame([test_data])

    def read_combined_data(self):
        """Read and combine data from both pagi and siang Excel files."""
        combined_data = []
        
        # Read pagi data
        if os.path.exists('wisuda_pagi.xlsx'):
            print("Reading wisuda_pagi.xlsx...")
            df_pagi = self.read_excel_data('wisuda_pagi.xlsx')
            if df_pagi is not None:
                df_pagi['SESI'] = 'Pagi'
                combined_data.append(df_pagi)
                print(f"  Found {len(df_pagi)} students in Pagi session")
        else:
            print("Warning: wisuda_pagi.xlsx not found")
        
        # Read siang data
        if os.path.exists('wisuda_siang.xlsx'):
            print("Reading wisuda_siang.xlsx...")
            df_siang = self.read_excel_data('wisuda_siang.xlsx')
            if df_siang is not None:
                df_siang['SESI'] = 'Siang'
                combined_data.append(df_siang)
                print(f"  Found {len(df_siang)} students in Siang session")
        else:
            print("Warning: wisuda_siang.xlsx not found")
        
        if not combined_data:
            print("Error: No data found in both Excel files")
            return None
        
        # Combine all data
        df_combined = pd.concat(combined_data, ignore_index=True)
        print(f"\nTotal combined data: {len(df_combined)} students")
        return df_combined

    def process_graduation_data(self, output_dir='output_revisi_pt', test_mode=False):
        """Main function to process graduation data."""
        if test_mode:
            print("=== TEST MODE: Generating single PPT with random data ===")
            df = self.create_test_data()
            print("Using test data for textbox position testing")
        else:
            print("Processing graduation data from Excel files...")
            df = self.read_combined_data()
            if df is None:
                return

            required_columns = [
                'PROGRAM STUDI', 'NAMA MAHASISWA', 'NIM', 'IPK', 'SKOR TAK',
                'Nama Dosen Wali', 'Nama Dosen Pembimbing 1', 'Nama Dosen Pembimbing 2',
                'PREDIKAT KELULUSAN', 'TEMPAT DUDUK', 'SESI'
            ]
            
            missing = [c for c in required_columns if c not in df.columns]
            if missing:
                print(f"Warning: Missing columns: {missing}")
                print(f"Available columns: {list(df.columns)}")

            if 'PREDIKAT KELULUSAN' in df.columns:
                print("\nPredikat distribution:")
                counts = df['PREDIKAT KELULUSAN'].value_counts()
                for p, c in counts.items():
                    t = self.get_predikat_template(p)
                    print(f"  {p}: {c} students -> {t} template")

        self.generate_ppt_revisi(df, output_dir, test_mode)
            
        if test_mode:
            print(f"\nTest PPT generated! Check the '{output_dir}' folder for 'TEST_POSITION.pptx'")
        else:
            print(f"\nProcessing completed! Check the '{output_dir}' folder for generated PPT files.")

def load_config():
    """Load configuration from config.json file."""
    config_file = 'config.json'
    default_config = {"TEST_MODE": False}
    
    try:
        if os.path.exists(config_file):
            with open(config_file, 'r', encoding='utf-8') as f:
                config = json.load(f)
                return config
        else:
            # Create default config file if it doesn't exist
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(default_config, f, indent=4)
                print(f"Created default config file: {config_file}")
            return default_config
    except Exception as e:
        print(f"Error loading config: {e}. Using default settings.")
        return default_config

def main():
    generator = GraduationPPTGenerator()

    # Load TEST_MODE from config file
    config = load_config()
    TEST_MODE = config.get('TEST_MODE', False)

    if TEST_MODE:
        print(f"\n{'='*50}")
        print("TEST MODE: Generating single PPT for textbox position testing")
        print(f"{'='*50}")
        output_dir = os.path.join('output_revisi_pt', 'Test')
        generator.process_graduation_data(output_dir, test_mode=True)
    else:
        # Process combined data with session separation
        print(f"\n{'='*50}")
        print("Processing data with session and seat separation")
        print(f"{'='*50}")
        
        output_dir = 'output_revisi_pt'
        generator.process_graduation_data(output_dir)

if __name__ == "__main__":
    main()
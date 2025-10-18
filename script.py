import os
import re
import pandas as pd
import random
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class GraduationPPTGenerator:
    DPI = 96  # konsisten dengan PowerPoint
    
    # POSISI FOTO - TENGAH FRAME MERAH
    # Sesuaikan nilai ini agar foto berada di tengah frame merah
    PHOTO_FRAME_W_CM = 5.0     # Lebar frame foto
    PHOTO_FRAME_H_CM = 7.0      # Tinggi frame foto
    FRAME_LEFT_CM = 7.0         # Posisi horizontal (tengah slide)
    FRAME_TOP_CM = 4.5          # Posisi vertikal (tengah frame merah)

    def __init__(self):
        self.templates = {
            'CUMLAUDE': 'templates/bg_cumlaude.png',
            'SUMMA CUMLAUDE': 'templates/bg_summa.png',
            'Non Predikat': 'templates/bg_non_predikat.png'
        }

    # =========================
    # Helpers ukuran & gambar
    # =========================
    def _set_slide_size_to_image_exact(self, prs, image_path, dpi=None):
        """Sesuaikan ukuran slide PERSIS dengan ukuran gambar (pixel -> inch @DPI)."""
        if dpi is None:
            dpi = self.DPI
        if not image_path or not os.path.exists(image_path):
            return
        with Image.open(image_path) as img:
            w_px, h_px = img.size
        prs.slide_width  = Inches(w_px / dpi)
        prs.slide_height = Inches(h_px / dpi)

    def _set_background_image(self, slide, image_path):
        """Pasang gambar sebagai latar dengan menambahkan picture di (0,0) ukuran native."""
        if not image_path or not os.path.exists(image_path):
            return
        try:
            slide.shapes.add_picture(image_path, 0, 0)
        except Exception as e:
            print(f"Error setting background image {image_path}: {e}")

    def _add_picture_fit(self, slide, image_path, left, top, frame_width, frame_height):
        """Tambahkan gambar agar pas di dalam frame tanpa distorsi (centered)."""
        try:
            with Image.open(image_path) as img:
                img_w, img_h = img.size
        except Exception as e:
            print(f"Error opening image {image_path}: {e}")
            return None

        frame_ratio = float(frame_width) / float(frame_height) if frame_height else 1.0
        img_ratio = img_w / img_h if img_h else 1.0

        if img_ratio > frame_ratio:
            width = int(frame_width)
            height = int(width / img_ratio)
        else:
            height = int(frame_height)
            width = int(height * img_ratio)

        int_left = int(left)
        int_top = int(top)
        int_fw = int(frame_width)
        int_fh = int(frame_height)
        offset_left = int_left + (int_fw - int(width)) // 2
        offset_top = int_top + (int_fh - int(height)) // 2

        try:
            return slide.shapes.add_picture(image_path, int(offset_left), int(offset_top), width=int(width), height=int(height))
        except Exception as e:
            print(f"Error adding fitted picture {image_path}: {e}")
            return None

    # =========================
    # Data helpers
    # =========================
    def read_excel_data(self, file_path):
        """Read Excel file and return DataFrame."""
        try:
            df = pd.read_excel(file_path)
            print(f"Successfully read {len(df)} rows from {file_path}")
            print(f"Columns: {list(df.columns)}")
            return df
        except Exception as e:
            print(f"Error reading Excel file: {e}")
            return None

    def get_predikat_template(self, predikat):
        """Determine template based on predikat kelulusan."""
        if pd.isna(predikat) or predikat == '':
            return 'Non Predikat'
        predikat_str = str(predikat).lower().strip()
        if 'summa' in predikat_str and 'cumlaude' in predikat_str:
            return 'SUMMA CUMLAUDE'
        elif 'cumlaude' in predikat_str and 'summa' not in predikat_str:
            return 'CUMLAUDE'
        else:
            return 'Non Predikat'

    def find_student_photo(self, nim, program_folder):
        """Find student photo based on NIM in program folder."""
        photo_path = os.path.join('photos', program_folder, f"{nim}_graduation_1.jpg")
        return photo_path if os.path.exists(photo_path) else None

    def extract_seat_position(self, tempat_duduk):
        """Extract seat position for ordering (format '1.1.L')."""
        if pd.isna(tempat_duduk) or str(tempat_duduk).strip() == '':
            return (999, 999, 'Z')
        try:
            parts = str(tempat_duduk).split('.')
            if len(parts) >= 3:
                row = int(parts[0]); seat = int(parts[1]); side = parts[2].upper()
                return (row, seat, side)
        except Exception:
            pass
        return (999, 999, 'Z')

    # =========================
    # Slide builders
    # =========================
    def create_slide(self, prs, student_data, photo_path):
        """Create a single slide for a student."""
        predikat = self.get_predikat_template(student_data.get('PREDIKAT KELULUSAN', ''))
        template_path = self.templates.get(predikat, self.templates['Non Predikat'])

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background full-bleed
        if os.path.exists(template_path):
            self._set_background_image(slide, template_path)

        # Posisi dan ukuran frame foto dalam CM
        frame_left, frame_top = Cm(self.FRAME_LEFT_CM), Cm(self.FRAME_TOP_CM)
        frame_w, frame_h = Cm(self.PHOTO_FRAME_W_CM), Cm(self.PHOTO_FRAME_H_CM)

        # FOTO: fit ke dalam frame merah (tengah)
        if photo_path and os.path.exists(photo_path):
            try:
                self._add_picture_fit(slide, photo_path, frame_left, frame_top, frame_w, frame_h)
            except Exception as e:
                print(f"Error adding photo {photo_path}: {e}")

        # Teks info mahasiswa & dosen
        self.add_student_info(slide, student_data)

        return slide

    def _add_textbox(self, slide, text, left, top, width, height, font_size=18, bold=False, upper=True, alignment=PP_ALIGN.LEFT):
        """Utility untuk menambah textbox konsisten."""
        if not text or str(text).strip().lower() == 'nan':
            return
        content = str(text).upper() if upper else str(text)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = content
        p.alignment = alignment

        if len(p.runs) == 0:
            run = p.add_run()
            run.text = content
            p.text = ''

        for run in p.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

    def add_student_info(self, slide, student_data):
        """Add student information to slide - POSISI BARU SESUAI TEMPLATE."""
        program = student_data.get('PROGRAM STUDI', '')
        nama = student_data.get('NAMA MAHASISWA', '')
        nim = student_data.get('NIM', '')
        ipk = student_data.get('IPK', '')
        tak = student_data.get('SKOR TAK', '')
        dosen_wali = student_data.get('Nama Dosen Wali', '')

        dosen_pembimbing1 = student_data.get('Nama Dosen Pembimbing 1', '')
        dosen_pembimbing2 = student_data.get('Nama Dosen Pembimbing 2', '')
        pembimbing_names = []
        if pd.notna(dosen_pembimbing1) and str(dosen_pembimbing1).strip() != '':
            pembimbing_names.append(str(dosen_pembimbing1))
        if pd.notna(dosen_pembimbing2) and str(dosen_pembimbing2).strip() != '':
            pembimbing_names.append(str(dosen_pembimbing2))

        # ==========================================
        # POSISI BARU SESUAI TEMPLATE
        # ==========================================
        
        # PROGRAM STUDI : 
        self._add_textbox(slide, program, Cm(4.5), Cm(3), Cm(10), Cm(1), font_size=14, bold=True, alignment=PP_ALIGN.CENTER)

        # # NAMA MAHASISWA : 
        self._add_textbox(slide, nama, Cm(0.2), Cm(14), Cm(19), Cm(1), font_size=19, bold=True, alignment=PP_ALIGN.CENTER)

        # # NIM : 
        self._add_textbox(slide, nim, Cm(4.4), Cm(15.36), Cm(4), Cm(0.8), font_size=16, bold=True)
        
        # # IPK : 
        self._add_textbox(slide, ipk, Cm(11.85), Cm(15.35), Cm(2), Cm(0.8), font_size=16, bold=True)
        
        # # TAK : 
        self._add_textbox(slide, tak, Cm(15.2), Cm(15.35), Cm(2), Cm(0.8), font_size=16, bold=True)
        
        # # DOSEN WALI : 
        self._add_textbox(slide, dosen_wali, Cm(6.8), Cm(16.3), Cm(12), Cm(0.8), font_size=14, bold=True)
        
        # DOSEN PEMBIMBING : 
        # Dosen Pembimbing Names
        if len(pembimbing_names) > 0:
            names_box = slide.shapes.add_textbox(Cm(6.8), Cm(17.1), Cm(12), Cm(1.5))
            names_tf = names_box.text_frame
            names_tf.clear()

            # Baris pertama
            p_name = names_tf.paragraphs[0]
            p_name.alignment = PP_ALIGN.LEFT
            run_n1 = p_name.add_run()
            run_n1.text = str(pembimbing_names[0]).upper()
            run_n1.font.name = 'Arial'
            run_n1.font.size = Pt(14)
            run_n1.font.bold = True
            run_n1.font.color.rgb = RGBColor(0, 0, 0)

            # Baris berikutnya
            for extra in pembimbing_names[1:]:
                p_more = names_tf.add_paragraph()
                p_more.alignment = PP_ALIGN.LEFT
                run_more = p_more.add_run()
                run_more.text = str(extra).upper()
                run_more.font.name = 'Arial'
                run_more.font.size = Pt(14)
                run_more.font.bold = True
                run_more.font.color.rgb = RGBColor(0, 0, 0)

    # =========================
    # Pipeline
    # =========================
    def generate_ppt_per_program(self, df, output_dir='output', test_mode=False):
        """Generate separate PPT files for each program."""
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if test_mode:
            # Test mode: generate single PPT with test data
            print(f"\nProcessing test data...")
            program_data = df.copy()
            
            prs = Presentation()
            
            # Use Cumlaude template for testing
            template_path = self.templates['CUMLAUDE']
            self._set_slide_size_to_image_exact(prs, template_path)
            
            # Add single test slide
            student = program_data.iloc[0]
            nim = student.get('NIM', '')
            photo_path = self.find_student_photo(nim, 'S1 Teknik Informatika')
            if photo_path:
                print(f"  Adding test slide for {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
            else:
                print(f"  Warning: Photo not found for test data (NIM: {nim})")
            self.create_slide(prs, student, photo_path)
            
            # Save test file
            output_file = os.path.join(output_dir, "TEST_POSITION.pptx")
            prs.save(output_file)
            print(f"  Saved: {output_file} (1 test slide)")
            return

        programs = df['PROGRAM STUDI'].unique()

        for program in programs:
            if pd.isna(program) or str(program).strip() == '':
                continue

            print(f"\nProcessing program: {program}")
            program_data = df[df['PROGRAM STUDI'] == program].copy()

            # Sort by seat position
            program_data['seat_sort'] = program_data['TEMPAT DUDUK'].apply(self.extract_seat_position)
            program_data = program_data.sort_values('seat_sort')

            prs = Presentation()

            # Tentukan template agar ukuran slide match EXACT background
            try:
                first_row = program_data.iloc[0]
                first_template_name = self.get_predikat_template(first_row.get('PREDIKAT KELULUSAN', ''))
                first_template_path = self.templates.get(first_template_name, self.templates['Non Predikat'])
            except Exception:
                first_template_path = self.templates['Non Predikat']

            self._set_slide_size_to_image_exact(prs, first_template_path)

            # Tambahkan slide per mahasiswa
            for _, student in program_data.iterrows():
                nim = student.get('NIM', '')
                photo_path = self.find_student_photo(nim, program)
                if photo_path:
                    print(f"  Adding slide for {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
                else:
                    print(f"  Warning: Photo not found for {student.get('NAMA MAHASISWA', '')} (NIM: {nim})")
                self.create_slide(prs, student, photo_path)

            # Simpan file
            safe_program_name = re.sub(r'[^\w\s-]', '', str(program)).strip()
            safe_program_name = re.sub(r'[-\s]+', '_', safe_program_name)
            output_file = os.path.join(output_dir, f"{safe_program_name}.pptx")
            prs.save(output_file)
            print(f"  Saved: {output_file} ({len(program_data)} slides)")

    def create_test_data(self):
        """Create random test data for testing textbox positions."""
        test_data = {
            'PROGRAM STUDI': 'S1 Rekayasa Perangkat Lunak',
            'NAMA MAHASISWA': 'JOHN DOE TESTING DAN TESTING',
            'NIM': '1201200001',
            'IPK': '3.85',
            'SKOR TAK': '450',
            'Nama Dosen Wali': 'Dr. Ahmad Wijaya, S.T., M.T.',
            'Nama Dosen Pembimbing 1': 'Prof. Dr. Budi Santoso, S.T., M.T.',
            'Nama Dosen Pembimbing 2': 'Dr. Citra Dewi, S.T., M.Kom.',
            'PREDIKAT KELULUSAN': 'Cumlaude',
            'TEMPAT DUDUK': '1.1.L'
        }
        return pd.DataFrame([test_data])

    def process_graduation_data(self, excel_file, output_dir='output', test_mode=False):
        """Main function to process graduation data."""
        if test_mode:
            print("=== TEST MODE: Generating single PPT with random data ===")
            df = self.create_test_data()
            print("Using test data for textbox position testing")
        else:
            print(f"Processing graduation data from: {excel_file}")
            df = self.read_excel_data(excel_file)
            if df is None:
                return

            required_columns = [
                'PROGRAM STUDI', 'NAMA MAHASISWA', 'NIM', 'IPK', 'SKOR TAK',
                'Nama Dosen Wali', 'Nama Dosen Pembimbing 1', 'Nama Dosen Pembimbing 2',
                'PREDIKAT KELULUSAN', 'TEMPAT DUDUK'
            ]
            missing = [c for c in required_columns if c not in df.columns]
            if missing:
                print(f"Warning: Missing columns: {missing}")
                print(f"Available columns: {list(df.columns)}")

            if 'PREDIKAT KELULUSAN' in df.columns:
                print("\nPredikat distribution:")
                counts = df['PREDIKAT KELULUSAN'].value_counts()
                for p, c in counts.items():
                    t = self.get_predikat_template(p)
                    print(f"  {p}: {c} students -> {t} template")

        self.generate_ppt_per_program(df, output_dir, test_mode)
        if test_mode:
            print(f"\nTest PPT generated! Check the '{output_dir}' folder for 'TEST_POSITION.pptx'")
        else:
            print(f"\nProcessing completed! Check the '{output_dir}' folder for generated PPT files.")

def main():
    generator = GraduationPPTGenerator()

    # TEST MODE: Set to True untuk testing posisi textbox
    TEST_MODE = True  # Ubah ke True untuk testing

    if TEST_MODE:
        print(f"\n{'='*50}")
        print("TEST MODE: Generating single PPT for textbox position testing")
        print(f"{'='*50}")
        output_dir = os.path.join('output', 'Test')
        generator.process_graduation_data('', output_dir, test_mode=True)
    else:
        # Daftar file excel & subfolder output
        excel_files = [
            ('wisuda_pagi.xlsx', 'Wisuda Pagi'),
            ('wisuda_siang.xlsx', 'Wisuda Siang')
        ]

        for excel_file, folder_name in excel_files:
            if os.path.exists(excel_file):
                print(f"\n{'='*50}")
                print(f"Processing: {excel_file}")
                print(f"Output folder: {folder_name}")
                print(f"{'='*50}")
                output_dir = os.path.join('output', folder_name)
                generator.process_graduation_data(excel_file, output_dir)
            else:
                print(f"File not found: {excel_file}")

if __name__ == "__main__":
    main()